import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
import textwrap

# Função para exibir o menu e retornar a escolha
def menu():
    menu_text = """
    Qual tipo de gráfico você deseja apresentar?
    [1] - Radar
    [2] - Colunas
    [3] - Ambos
    => """
    return int(input(textwrap.dedent(menu_text)))

# Lê os dados do Excel
df = pd.read_excel("dados.xlsx")

# Função para escolher qual template abrir
def escolhe_grafico():
    escolha = menu()
    if escolha == 1:
        return ["template"]
    elif escolha == 2:
        return ["template2"]
    elif escolha == 3:
        return ["template", "template2"]
    else:
        print("❌ Escolha não disponível")
        exit()

# Retorna uma lista de templates (1 ou mais)
templates = escolhe_grafico()

for template in templates:
    prs = Presentation(f"{template}.pptx")

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_chart:
                chart = shape.chart

                chart_data = CategoryChartData()
                chart_data.categories = df.iloc[:, 0].tolist()
                chart_data.add_series("Série 1", df.iloc[:, 1].tolist())

                chart.replace_data(chart_data)

    # Nome do arquivo final baseado no template
    prs.save(f"relatorio_final_{template}.pptx")
    print(f"✅ Arquivo gerado: relatorio_final_{template}.pptx")

print("🎯 Finalizado com sucesso!")
